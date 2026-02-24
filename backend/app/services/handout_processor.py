"""Handout processor: PDF/PPTX/VTT extraction, LLM parsing, chunking, and vectorization."""

from __future__ import annotations

import io
import re

import fitz  # PyMuPDF


def extract_text_from_pdf(pdf_bytes: bytes) -> str:
    """Extract text from PDF bytes using PyMuPDF."""
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n".join(pages)


def extract_text_from_pptx(pptx_bytes: bytes) -> str:
    """Extract text from PowerPoint (.pptx) bytes."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
        slides.append("\n".join(parts))
    return "\n\n".join(slides)


def extract_text_from_vtt(vtt_bytes: bytes) -> str:
    """Extract plain text from WebVTT (.vtt) subtitle files, stripping timestamps."""
    text = vtt_bytes.decode("utf-8", errors="replace")
    lines = text.split("\n")
    content_lines = []
    # Skip WEBVTT header, timestamps (HH:MM:SS.mmm --> HH:MM:SS.mmm), sequence numbers, and blank lines
    timestamp_re = re.compile(r"\d{2}:\d{2}[:\.].*-->")
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith("WEBVTT") or line.startswith("NOTE"):
            continue
        if timestamp_re.match(line):
            continue
        if line.isdigit():
            continue
        # Strip HTML-like tags (<b>, <i>, etc.)
        clean = re.sub(r"<[^>]+>", "", line)
        if clean:
            content_lines.append(clean)
    return " ".join(content_lines)


def extract_text(filename: str, file_bytes: bytes) -> str:
    """Extract text from a file based on its extension."""
    name = filename.lower()
    if name.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    elif name.endswith(".pptx"):
        return extract_text_from_pptx(file_bytes)
    elif name.endswith(".vtt"):
        return extract_text_from_vtt(file_bytes)
    else:
        # Plain text (.txt, .md, etc.)
        return file_bytes.decode("utf-8")


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    """Split text into word-based chunks of approximately *chunk_size* words
    with *overlap* word overlap between consecutive chunks."""
    words = text.split()
    if len(words) <= chunk_size:
        return [text]

    chunks: list[str] = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk_words = words[start:end]
        chunks.append(" ".join(chunk_words))
        if end >= len(words):
            break
        start = end - overlap
    return chunks


PARSE_PROMPT = """You are a course content analyzer. Analyze the following course handout and extract a complete course structure.

Return ONLY valid JSON (no markdown fences) in this exact format:
{
  "name": "Course Name (inferred from the handout content)",
  "description": "A 1-2 sentence description of what this course covers",
  "sections": [
    {
      "title": "Section Title",
      "summary": "Brief section summary",
      "learning_objectives": ["Objective 1", "Objective 2"],
      "key_concepts": ["Concept 1", "Concept 2"],
      "prerequisites": ["Prerequisite 1"],
      "subtopics": [
        {
          "title": "Subtopic Title",
          "content": "Full subtopic content from the handout",
          "summary": "Brief subtopic summary"
        }
      ]
    }
  ]
}

Rules:
- Infer the course name and description from the handout content
- Group related content logically into sections
- Each subtopic should be a focused, teachable unit
- Learning objectives should describe what students will be able to do after studying the section
- Key concepts are the important terms and ideas in the section
- Prerequisites are what students should know before this section

HANDOUT:
"""


class HandoutProcessor:
    """Parse handout text into sections/subtopics, chunk, embed, and index."""

    def __init__(self, db, llm, embedder):
        self.db = db
        self.llm = llm
        self.embedder = embedder

    async def process(self, course_id: str, text: str):
        """Parse handout text into sections/subtopics, chunk, embed, and index.

        Returns a dict with ``name`` and ``description`` auto-extracted from
        the handout content.
        """
        from app.services.vector_store import _to_blob

        result = await self.llm.chat_json(
            [{"role": "user", "content": PARSE_PROMPT + text}]
        )

        name = result.get("name", "")
        description = result.get("description", "")
        sections = result.get("sections", [])

        # Clear old sections before creating new ones (handles reprocessing)
        await self.db.delete_sections_by_course(course_id)

        for sec_idx, section in enumerate(sections):
            section_id = await self.db.create_section(
                course_id,
                section["title"],
                section.get("summary", ""),
                sec_idx,
                learning_objectives=section.get("learning_objectives", []),
                key_concepts=section.get("key_concepts", []),
                prerequisites=section.get("prerequisites", []),
            )

            for st_idx, subtopic in enumerate(section.get("subtopics", [])):
                subtopic_id = await self.db.create_subtopic(
                    section_id,
                    subtopic["title"],
                    subtopic.get("content", ""),
                    subtopic.get("summary", ""),
                    st_idx,
                )

                content = subtopic.get("content", "")
                if not content:
                    continue
                chunks = chunk_text(content)

                # Embed chunks if embedder available, store inline with content
                embeddings = None
                if self.embedder:
                    embeddings = await self.embedder.embed_batch(chunks)

                for i, chunk_text_val in enumerate(chunks):
                    emb_blob = _to_blob(embeddings[i]) if embeddings else None
                    await self.db.create_chunk_with_embedding(
                        subtopic_id, chunk_text_val, emb_blob, i
                    )

        return {"name": name, "description": description}
